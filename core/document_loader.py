"""
文档加载器 —— Layout-aware + MultiModal 解析

PDF: PyMuPDF blocks + 表格检测 + 内嵌图片 OCR + 元素类型分类
图片: EasyOCR + 布局感知
Word/Excel/PPT: 纯文本提取
"""

import os
import io
import logging
from io import StringIO

logger = logging.getLogger(__name__)


# ── OCR ──

_ocr = None


def preload_ocr():
    """启动时预加载 EasyOCR 模型"""
    _get_ocr()
    logger.info("EasyOCR 模型已预加载")


def _get_ocr():
    global _ocr
    if _ocr is None:
        import easyocr
        _ocr = easyocr.Reader(['ch_sim', 'en'], gpu=True)
    return _ocr


def _ocr_image(image):
    import numpy as np
    ocr = _get_ocr()
    img_array = np.array(image.convert("RGB"))
    try:
        results = ocr.readtext(img_array)
        if not results:
            return ""
        return "\n".join(text for _, text, _ in results)
    except Exception as e:
        logger.warning(f"EasyOCR 失败: {e}")
        return ""


# ── 元素类型分类 ──

def _classify_element(text, block_type, bbox, page_height, font_sizes=None):
    """
    根据文本特征 + 位置信息分类元素类型。
    PyMuPDF dict 模式可提供 font_sizes。

    分类: Title / SectionHeader / Paragraph / ListItem / Table / Caption / HeaderFooter / Image
    """
    text_stripped = text.strip()

    # Header/Footer: 页面顶部 15% 或底部 15%
    if bbox and page_height:
        y0, y1 = bbox[1], bbox[3]
        if y1 < page_height * 0.15 or y0 > page_height * 0.85:
            return "HeaderFooter"

    # Image block
    if block_type == 1:
        return "Image"

    # Table 块 (由 find_tables 标记)
    if block_type == 2:
        return "Table"

    # Title: 短文本 + 大字体（如果没有 font_size 信息则检测全部大写或编号模式）
    if len(text_stripped) < 80:
        if font_sizes and max(font_sizes) > 14:
            return "Title"
        # Heuristic: 独立短行，不以标点结尾
        if text_stripped[-1] not in "，。；：！？,.;:!?" and len(text_stripped) > 2:
            # 检测编号模式: "第X章", "1.", "一、"
            if any(text_stripped.startswith(p) for p in
                   ["第", "CHAPTER", "Chapter"]):
                return "Title"
            if text_stripped[:2].isdigit() and text_stripped[2:3] in ".、．":
                return "Title"

    # Caption: "图X", "表X", "Figure", "Table" 开头
    if any(text_stripped.startswith(p) for p in
           ["图", "表", "Figure", "Table", "FIG"]):
        return "Caption"

    # ListItem: 以 -, *, •, 1., a) 等开头
    first_line = text_stripped.split("\n")[0].strip()
    list_prefixes = ["- ", "* ", "• ", "+ ", "· "]
    if any(first_line.startswith(p) for p in list_prefixes):
        return "ListItem"
    if len(first_line) > 2 and first_line[0].isdigit() and first_line[1] in ".、)）":
        return "ListItem"

    # 全部大写 → SectionHeader
    if text_stripped.isupper() and len(text_stripped) < 60:
        return "SectionHeader"

    return "Paragraph"


# ── PDF: PyMuPDF ──

def _extract_pdf_pymupdf(filepath):
    """
    PyMuPDF 提取 PDF：blocks + 表格 + 内嵌图片。
    表格优先使用 unstructured 的 HTML 表格输出，回退 PyMuPDF find_tables。
    返回 layout-aware blocks 列表。
    """
    import fitz
    doc = fitz.open(filepath)
    blocks = []

    # ── unstructured 表格提取（整份文档，page-by-page 映射）──
    unstructured_tables = {}  # page_num → [(bbox, html_table), ...]
    try:
        from unstructured.partition.pdf import partition_pdf
        elements = partition_pdf(
            filename=filepath, strategy="auto",
            infer_table_structure=True,
        )
        for el in elements:
            cat = el.category if hasattr(el, "category") else str(type(el))
            if "Table" in cat:
                meta = el.metadata.to_dict() if hasattr(el.metadata, "to_dict") else {}
                page = (meta.get("page_number", 1)
                        if isinstance(meta, dict) else 1)
                # unstructured 的 table 输出已是 HTML
                html = el.text if hasattr(el, "text") else str(el)
                # 简单 HTML → Markdown table
                text = _html_table_to_text(html)
                unstructured_tables.setdefault(page, []).append((text, meta))
    except Exception as e:
        logger.warning(f"unstructured 表格提取失败 ({e})，使用 PyMuPDF fallback")

    for page_num, page in enumerate(doc, 1):
        page_height = page.rect.height

        # 1. 优先使用 unstructured 表格
        page_tables = unstructured_tables.get(page_num, [])
        if page_tables:
            for table_text, meta in page_tables:
                blocks.append({
                    "text": table_text, "page": page_num,
                    "bbox": None, "block_type": 2, "element_type": "Table",
                    "source_type": "pdf_table",
                })
        else:
            # 回退 PyMuPDF 表格
            try:
                tables = page.find_tables()
                for table in tables:
                    header = table.header.names if table.header else []
                    rows = []
                    for row in table.extract():
                        rows.append(" | ".join(str(c) if c else "" for c in row))
                    table_text = f"[表格]\n" + "\n".join(rows)
                    blocks.append({
                        "text": table_text, "page": page_num,
                        "bbox": tuple(table.bbox) if table.bbox else None,
                        "block_type": 2, "element_type": "Table",
                        "source_type": "pdf",
                    })
            except Exception:
                pass

        # 2. 提取内嵌图片并 OCR
        try:
            for img_info in page.get_images(full=True):
                xref = img_info[0]
                base_image = doc.extract_image(xref)
                if base_image and base_image["width"] > 50 and base_image["height"] > 50:
                    img_bytes = base_image["image"]
                    try:
                        from PIL import Image
                        pil_img = Image.open(io.BytesIO(img_bytes))
                        ocr_text = _ocr_image(pil_img)
                        if ocr_text:
                            rects = page.get_image_rects(xref)
                            bbox = tuple(rects[0]) if rects else None
                            blocks.append({
                                "text": f"[图片描述]\n{ocr_text}", "page": page_num,
                                "bbox": bbox, "block_type": 1, "element_type": "Image",
                                "source_type": "pdf_image",
                            })
                    except Exception:
                        pass
        except Exception:
            pass

        # 3. 提取文本块（dict 模式获取字体）
        try:
            dict_blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        except Exception:
            dict_blocks = []

        for b in dict_blocks:
            if b.get("type") == 0:
                text = ""
                font_sizes = []
                for line in b.get("lines", []):
                    for span in line.get("spans", []):
                        text += span["text"]
                        font_sizes.append(span.get("size", 10))
                text = text.strip()
                if not text:
                    continue
                if text.isdigit() and len(text) <= 3:
                    continue
                bbox = tuple(b["bbox"])
                element_type = _classify_element(text, 0, bbox, page_height, font_sizes)
                if element_type == "HeaderFooter":
                    continue
                blocks.append({
                    "text": text, "page": page_num,
                    "bbox": bbox, "block_type": 0,
                    "element_type": element_type, "source_type": "pdf",
                })

    doc.close()
    return blocks


def _html_table_to_text(html):
    """将 unstructured 输出的 HTML 表格转为紧凑文本"""
    import re
    # 去除 HTML 标签，保留内容
    text = re.sub(r'</tr>', '\n', html)
    text = re.sub(r'</t[dh]>', ' | ', text)
    text = re.sub(r'<[^>]+>', '', text)
    text = re.sub(r'\n\s*\n', '\n', text)
    text = re.sub(r'\| $', '', text, flags=re.MULTILINE)
    return "[表格]\n" + text.strip()


def _extract_pdf_pdfminer(filepath):
    """pdfminer 纯文本回退方案"""
    from pdfminer.high_level import extract_text_to_fp
    output = StringIO()
    with open(filepath, 'rb') as f:
        extract_text_to_fp(f, output)
    return output.getvalue()


# ── 图片文件 ──

def _extract_image_file(filepath):
    """处理用户直接上传的图片文件，OCR 提取文字"""
    from PIL import Image
    img = Image.open(filepath)
    text = _ocr_image(img)
    if not text:
        raise ValueError("图片中未识别到文字内容")
    return [{"text": text, "page": 1, "bbox": None, "block_type": 0,
             "element_type": "Paragraph", "source_type": "image"}]


# ── 统一入口 ──

def extract_text(filepath):
    """
    多格式文档提取入口。

    Returns:
        layout-aware → [{"text","page","bbox","block_type","element_type","source_type"}, ...]
        纯文本       → str
    """
    file_ext = os.path.splitext(filepath)[1].lower()

    # PDF
    if file_ext == '.pdf':
        try:
            blocks = _extract_pdf_pymupdf(filepath)
            if not blocks:
                raise ValueError("PyMuPDF 未提取到内容")
            return blocks
        except Exception as e:
            logger.warning(f"PyMuPDF 提取失败 ({e})，回退 pdfminer")
            return _extract_pdf_pdfminer(filepath)

    # 图片
    if file_ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp', '.gif']:
        return _extract_image_file(filepath)

    # 纯文本 / Markdown
    if file_ext in ['.txt', '.md']:
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()

    # Word
    if file_ext == '.docx':
        try:
            from docx import Document
            doc = Document(filepath)
            return "\n".join(p.text for p in doc.paragraphs)
        except ImportError:
            logger.error("处理Word文档需要安装python-docx库")
            return ""

    # Excel
    if file_ext in ['.xlsx', '.xls']:
        try:
            import pandas as pd
            text_parts = []
            for sheet_name, df in pd.read_excel(filepath, sheet_name=None).items():
                text_parts.append(f"[工作表: {sheet_name}]\n{df.to_string(index=False)}")
            return "\n\n".join(text_parts)
        except ImportError:
            logger.error("处理Excel文件需要安装pandas库")
            return ""

    # PPT
    if file_ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            slides = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slides.append(shape.text.strip())
            return "\n".join(slides)
        except ImportError:
            logger.error("处理PPT文件需要安装python-pptx库")
            return ""

    logger.warning(f"不支持的文件格式: {file_ext}")
    return ""


def is_layout_aware_result(result):
    """判断提取结果是否为 layout-aware 格式"""
    return isinstance(result, list)


def blocks_to_plain_text(blocks):
    """将 layout-aware blocks 合并为纯文本"""
    return "\n\n".join(b["text"] for b in blocks)
